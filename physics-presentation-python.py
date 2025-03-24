from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

def create_physics_presentation():
    # プレゼンテーションの作成
    prs = Presentation()
    
    # スライドサイズを16:9に設定（通常のワイドスクリーン）
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # カラーの定義
    PRIMARY_COLOR = RGBColor(37, 99, 235)    # #2563EB ロイヤルブルー
    SECONDARY_COLOR = RGBColor(71, 85, 105)  # #475569 ダークグレー
    ACCENT_COLOR = RGBColor(248, 250, 252)   # #f8fafc ライトグレー
    
    RISK_HIGH = RGBColor(239, 68, 68)        # #ef4444 赤
    RISK_MEDIUM = RGBColor(245, 158, 11)     # #f59e0b オレンジ
    RISK_LOW = RGBColor(16, 185, 129)        # #10b981 緑
    
    # スライドレイアウトの選択と追加
    blank_slide_layout = prs.slide_layouts[6]  # 白紙のスライド
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # ====== ヘッダーの追加 ======
    header_box = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(0.5), 
        width=Inches(10), height=Inches(1)
    )
    header_text = header_box.text_frame
    header_text.text = "日常に潜む物理法則：5つの興味深い現象とその科学的解説"
    header_text.paragraphs[0].font.size = Pt(32)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # サブタイトル
    subtitle_p = header_text.add_paragraph()
    subtitle_p.text = "物理現象の専門的視点からの分析と応用可能性"
    subtitle_p.font.size = Pt(16)
    subtitle_p.font.color.rgb = SECONDARY_COLOR
    
    # 日付と出典
    date_box = slide.shapes.add_textbox(
        left=Inches(12), top=Inches(0.5), 
        width=Inches(3.5), height=Inches(0.5)
    )
    date_text = date_box.text_frame
    date_text.text = "2025年3月22日"
    date_text.paragraphs[0].font.size = Pt(12)
    date_text.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    date_text.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    source_p = date_text.add_paragraph()
    source_p.text = "物理学レポート"
    source_p.font.size = Pt(12)
    source_p.font.color.rgb = SECONDARY_COLOR
    source_p.alignment = PP_ALIGN.RIGHT

    # ====== メインコンテンツの左カラム (60%) ======
    left_column = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left=Inches(0.5), top=Inches(2), 
        width=Inches(9), height=Inches(6)
    )
    left_column.fill.solid()
    left_column.fill.fore_color.rgb = ACCENT_COLOR
    left_column.line.color.rgb = ACCENT_COLOR
    
    # エグゼクティブサマリーカード
    exec_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        left=Inches(0.7), top=Inches(2.2), 
        width=Inches(8.6), height=Inches(1.8)
    )
    exec_card.fill.solid()
    exec_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    exec_card.line.color.rgb = RGBColor(255, 255, 255)
    exec_card.shadow.inherit = False
    
    # エグゼクティブサマリータイトル
    exec_title = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(2.3), 
        width=Inches(8), height=Inches(0.3)
    )
    exec_title_frame = exec_title.text_frame
    exec_title_frame.text = "エグゼクティブサマリー"
    exec_title_frame.paragraphs[0].font.size = Pt(18)
    exec_title_frame.paragraphs[0].font.bold = True
    exec_title_frame.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)  # #1e40af
    
    # エグゼクティブサマリー本文
    exec_text = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(2.7), 
        width=Inches(8), height=Inches(0.6)
    )
    exec_text_frame = exec_text.text_frame
    exec_text_frame.text = "物理現象の数学的記述は単なる学術的関心を超え、産業イノベーションとビジネス競争力に直結している。"
    exec_text_frame.paragraphs[0].font.size = Pt(14)
    exec_text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 65, 85)  # #334155
    
    p2 = exec_text_frame.add_paragraph()
    p2.text = "本レポートでは5つの異なる複雑さを持つ現象を分析し、それらがもたらす潜在的なビジネスチャンスと工学的応用を検討する。"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(51, 65, 85)
    
    # KPI表示
    kpi_data = [
        {"value": "9.8 m/s²", "label": "地球の重力加速度"},
        {"value": "0.002〜0.006", "label": "転がり抵抗係数範囲"},
        {"value": "<0.001g", "label": "シャボン玉の質量"},
        {"value": "~1000", "label": "葉のレイノルズ数"},
        {"value": "180°", "label": "機体反転角度"}
    ]
    
    for i, kpi in enumerate(kpi_data):
        kpi_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            left=Inches(0.7 + i*1.72), top=Inches(3.5), 
            width=Inches(1.5), height=Inches(0.4)
        )
        kpi_box.fill.solid()
        kpi_box.fill.fore_color.rgb = RGBColor(241, 245, 249)  # #f1f5f9
        kpi_box.line.color.rgb = RGBColor(241, 245, 249)
        
        kpi_text = slide.shapes.add_textbox(
            left=Inches(0.7 + i*1.72), top=Inches(3.55), 
            width=Inches(1.5), height=Inches(0.3)
        )
        kpi_text_frame = kpi_text.text_frame
        kpi_text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        kpi_text_frame.word_wrap = True
        
        value_p = kpi_text_frame.paragraphs[0]
        value_p.text = kpi["value"]
        value_p.font.size = Pt(14)
        value_p.font.bold = True
        value_p.font.color.rgb = PRIMARY_COLOR
        value_p.alignment = PP_ALIGN.CENTER
        
        label_p = kpi_text_frame.add_paragraph()
        label_p.text = kpi["label"]
        label_p.font.size = Pt(10)
        label_p.font.color.rgb = SECONDARY_COLOR
        label_p.alignment = PP_ALIGN.CENTER

    # 比較表のカード
    comparison_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        left=Inches(0.7), top=Inches(4.1), 
        width=Inches(8.6), height=Inches(2.2)
    )
    comparison_card.fill.solid()
    comparison_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    comparison_card.line.color.rgb = RGBColor(255, 255, 255)
    
    # 比較表タイトル
    comparison_title = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(4.2), 
        width=Inches(8), height=Inches(0.3)
    )
    comparison_title_frame = comparison_title.text_frame
    comparison_title_frame.text = "複雑性と応用可能性の比較"
    comparison_title_frame.paragraphs[0].font.size = Pt(18)
    comparison_title_frame.paragraphs[0].font.bold = True
    comparison_title_frame.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)
    
    # 比較表ヘッダー
    table_headers = ["指標", "重力加速度", "タイヤ抵抗", "シャボン玉", "葉の落下", "背面飛行"]
    table_data = [
        ["数学的複雑性", "低", "中", "高", "高", "中"],
        ["日常での観察頻度", "高", "高", "中", "高", "低"],
        ["産業応用可能性", "高", "高", "低", "中", "高"]
    ]
    
    # テーブルの作成（簡易的な表現）
    table_top = Inches(4.6)
    for i, header in enumerate(table_headers):
        header_box = slide.shapes.add_textbox(
            left=Inches(1 + i*1.4), top=table_top, 
            width=Inches(1.3), height=Inches(0.3)
        )
        header_text = header_box.text_frame
        header_text.text = header
        header_text.paragraphs[0].font.size = Pt(12)
        header_text.paragraphs[0].font.bold = True
        header_text.paragraphs[0].font.color.rgb = SECONDARY_COLOR
        header_text.paragraphs[0].alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
    
    # テーブルの行データ
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_value in enumerate(row_data):
            cell_box = slide.shapes.add_textbox(
                left=Inches(1 + col_idx*1.4), top=table_top + Inches(0.4 + row_idx*0.3), 
                width=Inches(1.3), height=Inches(0.3)
            )
            cell_text = cell_box.text_frame
            cell_text.text = cell_value
            cell_text.paragraphs[0].font.size = Pt(12)
            cell_text.paragraphs[0].font.color.rgb = RGBColor(51, 65, 85)
            cell_text.paragraphs[0].alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            
            # 最初の列は太字
            if col_idx == 0:
                cell_text.paragraphs[0].font.bold = True

    # タイムライン分析カード
    timeline_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        left=Inches(0.7), top=Inches(6.4), 
        width=Inches(8.6), height=Inches(1.5)
    )
    timeline_card.fill.solid()
    timeline_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    timeline_card.line.color.rgb = RGBColor(255, 255, 255)
    
    # タイムラインタイトル
    timeline_title = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(6.5), 
        width=Inches(8), height=Inches(0.3)
    )
    timeline_title_frame = timeline_title.text_frame
    timeline_title_frame.text = "応用開発の市場展開予測"
    timeline_title_frame.paragraphs[0].font.size = Pt(18)
    timeline_title_frame.paragraphs[0].font.bold = True
    timeline_title_frame.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)
    
    # タイムラインデータ
    timeline_data = [
        {"period": "2025-2026", "description": "転がり抵抗の最適化による次世代モビリティの効率化（20%エネルギー削減）"},
        {"period": "2026-2028", "description": "落下運動の空気力学的特性を応用した新素材・デバイス設計フレームワークの確立"},
        {"period": "2028-2030", "description": "飛行力学の高度シミュレーションによる次世代航空機の設計革新と市場拡大"}
    ]
    
    for i, item in enumerate(timeline_data):
        # タイムライン期間
        period_box = slide.shapes.add_textbox(
            left=Inches(1.5), top=Inches(6.9 + i*0.3), 
            width=Inches(1.5), height=Inches(0.3)
        )
        period_text = period_box.text_frame
        period_text.text = item["period"]
        period_text.paragraphs[0].font.size = Pt(12)
        period_text.paragraphs[0].font.bold = True
        period_text.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # 説明テキスト
        desc_box = slide.shapes.add_textbox(
            left=Inches(3), top=Inches(6.9 + i*0.3), 
            width=Inches(6), height=Inches(0.3)
        )
        desc_text = desc_box.text_frame
        desc_text.text = item["description"]
        desc_text.paragraphs[0].font.size = Pt(12)
        desc_text.paragraphs[0].font.color.rgb = RGBColor(51, 65, 85)
        
        # タイムラインの丸いポイント
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            left=Inches(1.2), top=Inches(6.95 + i*0.3), 
            width=Inches(0.15), height=Inches(0.15)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = PRIMARY_COLOR
        dot.line.color.rgb = PRIMARY_COLOR

    # ====== メインコンテンツの右カラム (40%) ======
    right_column = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left=Inches(9.8), top=Inches(2), 
        width=Inches(5.7), height=Inches(6)
    )
    right_column.fill.solid()
    right_column.fill.fore_color.rgb = RGBColor(241, 245, 249)  # #f1f5f9
    right_column.line.color.rgb = RGBColor(241, 245, 249)
    
    # 物理現象データ
    physics_data = [
        {
            "title": "重力加速度と自由落下運動",
            "formula": "h = ½gt²",
            "risk_level": "低",
            "risk_color": RISK_LOW,
            "points": [
                "自由落下は万有引力の法則に基づき高さによって時間が決まる",
                "実際の落下速度は空気抵抗により減衰し終端速度に達する",
                "粘性抗力は軽量物体に特に影響を与える"
            ],
            "application": "精密測定機器、建築構造設計、宇宙工学"
        },
        {
            "title": "自転車のタイヤ幅と転がり抵抗",
            "formula": "Frr = Crr・N",
            "risk_level": "中",
            "risk_color": RISK_MEDIUM,
            "points": [
                "細いタイヤは接地面積が小さいため効率的な走行が可能",
                "転がり抵抗はタイヤの変形によるヒステリシス損失が主因",
                "ロードバイク(Crr≈0.002)はマウンテンバイク(Crr≈0.006)より効率的"
            ],
            "application": "次世代モビリティ、エネルギー効率最適化、スポーツ工学"
        },
        {
            "title": "シャボン玉の回転運動と空気力学",
            "formula": "非線形流体方程式群",
            "risk_level": "低",
            "risk_color": RISK_LOW,
            "points": [
                "クッタ条件や境界層分離点が回転運動に影響を与える",
                "非対称な流速分布が回転モーメントを生成する",
                "コアンダ効果とマグナス効果が複雑な運動を引き起こす"
            ],
            "application": "微小流体デバイス、医療用マイクロロボット、センサー技術"
        },
        {
            "title": "木の葉の落下と空気抵抗",
            "formula": "Re = ρvL/μ",
            "risk_level": "中",
            "risk_color": RISK_MEDIUM,
            "points": [
                "落下運動はレイノルズ数と物体形状に強く依存する",
                "小さなRe数と不規則な形状がカオス的振る舞いを生む",
                "上下振動や左右へのスラローム運動など複雑な軌道を描く"
            ],
            "application": "バイオミメティクス、空力設計、拡散モデル"
        },
        {
            "title": "航空機の背面飛行と揚力反転",
            "formula": "L = ½ρv²SC_L",
            "risk_level": "高",
            "risk_color": RISK_HIGH,
            "points": [
                "背面飛行では迎え角調整により逆向きの揚力を発生させる",
                "より大きな迎え角が必要なためエネルギー消費が増加する",
                "機体構造や重心位置の非対称性が安定性に影響を与える"
            ],
            "application": "航空宇宙技術、姿勢制御システム、自律飛行ドローン"
        }
    ]
    
    # 各物理現象のカードを作成
    card_height = Inches(1.05)  # カードの高さ
    
    for i, data in enumerate(physics_data):
        # カード本体
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            left=Inches(10), top=Inches(2.2 + i*1.15), 
            width=Inches(5.3), height=card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(255, 255, 255)
        
        # カードヘッダー
        card_header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            left=Inches(10), top=Inches(2.2 + i*1.15), 
            width=Inches(5.3), height=Inches(0.3)
        )
        card_header.fill.solid()
        card_header.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card_header.line.color.rgb = RGBColor(241, 245, 249)
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            left=Inches(10.2), top=Inches(2.2 + i*1.15), 
            width=Inches(4), height=Inches(0.3)
        )
        title_text = title_box.text_frame
        title_text.text = data["title"]
        title_text.paragraphs[0].font.size = Pt(14)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)
        
        # リスクレベル
        risk_pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            left=Inches(14.3), top=Inches(2.23 + i*1.15), 
            width=Inches(0.9), height=Inches(0.2)
        )
        risk_pill.fill.solid()
        risk_pill.fill.fore_color.rgb = data["risk_color"]
        risk_pill.line.color.rgb = data["risk_color"]
        
        risk_text = slide.shapes.add_textbox(
            left=Inches(14.3), top=Inches(2.23 + i*1.15), 
            width=Inches(0.9), height=Inches(0.2)
        )
        risk_tf = risk_text.text_frame
        risk_tf.text = f"複雑性: {data['risk_level']}"
        risk_tf.paragraphs[0].font.size = Pt(9)
        risk_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        risk_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 数式
        formula_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            left=Inches(10.2), top=Inches(2.55 + i*1.15), 
            width=Inches(5), height=Inches(0.25)
        )
        formula_box.fill.solid()
        formula_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
        formula_box.line.color.rgb = RGBColor(248, 250, 252)
        
        formula_text = slide.shapes.add_textbox(
            left=Inches(10.2), top=Inches(2.55 + i*1.15), 
            width=Inches(5), height=Inches(0.25)
        )
        formula_tf = formula_text.text_frame
        formula_tf.text = data["formula"]
        formula_tf.paragraphs[0].font.size = Pt(12)
        formula_tf.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        formula_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # ポイント
        for j, point in enumerate(data["points"]):
            point_text = slide.shapes.add_textbox(
                left=Inches(10.4), top=Inches(2.85 + i*1.15 + j*0.18), 
                width=Inches(4.8), height=Inches(0.18)
            )
            point_tf = point_text.text_frame
            point_tf.text = f"• {point}"
            point_tf.paragraphs[0].font.size = Pt(10)
            point_tf.paragraphs[0].font.color.rgb = RGBColor(51, 65, 85)

    # フッター
    footer_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left=Inches(0), top=Inches(8.2), 
        width=Inches(16), height=Inches(0.8)
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    footer_shape.line.color.rgb = RGBColor(255, 255, 255)
    
    # ロゴ
    logo_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        left=Inches(0.5), top=Inches(8.4), 
        width=Inches(0.3), height=Inches(0.3)
    )
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = PRIMARY_COLOR
    logo_box.line.color.rgb = PRIMARY_COLOR
    
    logo_text = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(8.4), 
        width=Inches(0.3), height=Inches(0.3)
    )
    logo_tf = logo_text.text_frame
    logo_tf.text = "P"
    logo_tf.paragraphs[0].font.size = Pt(14)
    logo_tf.paragraphs[0].font.bold = True
    logo_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    logo_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # フッターテキスト
    footer_text_left = slide.shapes.add_textbox(
        left=Inches(0.9), top=Inches(8.45), 
        width=Inches(5), height=Inches(0.3)
    )
    footer_text_left_tf = footer_text_left.text_frame
    footer_text_left_tf.text = "Physical Sciences Advanced Research Group"
    footer_text_left_tf.paragraphs[0].font.size = Pt(10)
    footer_text_left_tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    
    footer_text_right = slide.shapes.add_textbox(
        left=Inches(10), top=Inches(8.45), 
        width=Inches(5.5), height=Inches(0.3)
    )
    footer_text_right_tf = footer_text_right.text_frame
    footer_text_right_tf.text = "出典: 物理現象の専門的分析 | ID: PSC-2025-03-22"
    footer_text_right_tf.paragraphs[0].font.size = Pt(10)
    footer_text_right_tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    footer_text_right_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # プレゼンテーションの保存
    prs.save('physics_presentation.pptx')
    return "physics_presentation.pptx"

if __name__ == "__main__":
    output_file = create_physics_presentation()
    print(f"プレゼンテーションが {output_file} として保存されました。")
