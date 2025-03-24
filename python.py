# make_ppt_with_image.py

from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO

def main():
    # 画像URL
    img_url = "https://cdn-contents.anymindgroup.com/corporate/wp-uploads/2024/11/22092837/image7-1.png"

    # PowerPointプレゼンテーション作成
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙スライド

    # 画像をダウンロードしてメモリに保持
    response = requests.get(img_url)
    if response.status_code != 200:
        print("画像の取得に失敗しました")
        return

    image_stream = BytesIO(response.content)

    # スライドに画像を貼り付け
    left = Inches(1)    # x位置
    top = Inches(1)     # y位置
    width = Inches(5)   # 幅
    height = Inches(3)  # 高さ

    slide.shapes.add_picture(image_stream, left, top, width, height)

    # 保存
    prs.save("output.pptx")
    print("✅ PowerPointファイル output.pptx を作成しました！")

if __name__ == "__main__":
    main()
