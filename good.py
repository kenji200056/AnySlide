import xmltodict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re
import webcolors

# ✅ SVGファイルのパス
svg_file_path = "1.svg"

# ✅ PowerPointのスライドサイズ (16:9)
ppt_width = Inches(10)
ppt_height = Inches(5.625)

# ✅ デバッグ用フラグ
DEBUG = True

# ✅ SVGファイルの読み込み
with open(svg_file_path, "r", encoding="utf-8") as f:
    svg_data = f.read()
print("✅ SVG ファイルを読み込みました")

# ✅ SVGの解析
svg_json = xmltodict.parse(svg_data)

# ✅ viewBox の取得
viewBox = svg_json["svg"].get("@viewBox", "0 0 1600 900").split()
vbX, vbY, vbWidth, vbHeight = map(float, viewBox)
print(f"✅ viewBox: {vbX} {vbY} {vbWidth} {vbHeight}")

# ✅ スケール変換係数 (SVG → PPT)
scale_x = ppt_width / vbWidth
scale_y = ppt_height / vbHeight
scale = min(scale_x, scale_y)  # **等比スケールを適用**

# ✅ 色変換関数
def convert_color(color):
    if not color:
        return RGBColor(0, 0, 0)  # デフォルト黒
    if color.startswith("#"):
        return RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))
    try:
        rgb = webcolors.name_to_rgb(color)
        return RGBColor(rgb.red, rgb.green, rgb.blue)
    except ValueError:
        return RGBColor(37, 99, 235)  # **デフォルトは青 (#2563EB)**

# ✅ PowerPointの作成
prs = Presentation()
prs.slide_width = ppt_width
prs.slide_height = ppt_height
slide = prs.slides.add_slide(prs.slide_layouts[6])  # **白紙スライド**

# ✅ SVG要素をPPTXに変換する関数
def process_element(element, parent_transform=(0, 0)):
    if not isinstance(element, dict):
        return

    for tag, content in element.items():
        if not isinstance(content, dict):
            continue  # **無効な要素をスキップ**

        attrib = content.get("@", {})

        if DEBUG:
            print(f"🔹 検出: {tag} (属性: {attrib})")

        # **TEXT要素の処理**
        if tag == "text":
            text = content.get("#text", "").strip()
            if text:
                x = (float(attrib.get("@x", 0)) - vbX + parent_transform[0]) * scale
                y = (float(attrib.get("@y", 0)) - vbY + parent_transform[1]) * scale
                font_size = float(attrib.get("@font-size", "14").replace("px", "")) * 0.75
                text_color = convert_color(attrib.get("@fill", "#000000"))
                text_anchor = attrib.get("@text-anchor", "start")

                print(f"📌 TEXT 要素: \"{text}\"")
                print(f"➡ 座標: ({x:.2f}, {y:.2f}), サイズ: {font_size:.2f}pt, 色: {attrib.get('@fill', '#000000')}")

                textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(0.5))
                tf = textbox.text_frame
                tf.word_wrap = False
                p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(font_size)
                p.font.color.rgb = text_color

                if text_anchor == "middle":
                    textbox.left = Inches(x - 2)
                elif text_anchor == "end":
                    textbox.left = Inches(x - 4)

                print(f"✅ PowerPoint に追加: \"{text}\"")

        # **RECT要素の処理**
        elif tag == "rect":
            x = (float(attrib.get("@x", 0)) - vbX + parent_transform[0]) * scale
            y = (float(attrib.get("@y", 0)) - vbY + parent_transform[1]) * scale
            width = float(attrib.get("@width", "100")) * scale
            height = float(attrib.get("@height", "30")) * scale
            fill = convert_color(attrib.get("@fill", "#FFFFFF"))

            slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height)).fill.fore_color.rgb = fill
            print(f"✅ RECT追加: x={x:.2f}, y={y:.2f}, w={width:.2f}, h={height:.2f}, 色={attrib.get('@fill', '#FFFFFF')}")

        # **CIRCLE要素の処理**
        elif tag == "circle":
            cx = (float(attrib.get("@cx", 0)) - vbX + parent_transform[0]) * scale
            cy = (float(attrib.get("@cy", 0)) - vbY + parent_transform[1]) * scale
            r = float(attrib.get("@r", "10")) * scale
            fill = convert_color(attrib.get("@fill", "#000000"))

            slide.shapes.add_shape(1, Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2)).fill.fore_color.rgb = fill
            print(f"✅ CIRCLE追加: cx={cx:.2f}, cy={cy:.2f}, r={r:.2f}, 色={attrib.get('@fill', '#000000')}")

        # **G要素の処理（グループ要素）**
        elif tag == "g":
            new_transform = (float(attrib.get("@x", 0)) + parent_transform[0], float(attrib.get("@y", 0)) + parent_transform[1])
            process_element(content, new_transform)

        # **子要素の処理**
        process_element(content, parent_transform)

# ✅ SVG要素の変換
svg_elements = svg_json["svg"]
process_element(svg_elements)

# ✅ PPTXの保存
prs.save("presentation_fixed.pptx")
print("✅ PPTXファイルが作成されました: presentation_fixed.pptx")
